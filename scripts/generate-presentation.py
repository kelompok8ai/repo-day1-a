#!/usr/bin/env python3
"""Generate CorpSec Bank Sumut presentation (PPTX)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# Bank Sumut brand colors
NAVY = RGBColor(0x0C, 0x23, 0x40)
NAVY_MID = RGBColor(0x15, 0x34, 0x5C)
ORANGE = RGBColor(0xF5, 0x82, 0x20)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF4, 0xF6, 0xF9)
SLATE = RGBColor(0x64, 0x74, 0x8B)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)

OUTPUT = Path(__file__).resolve().parent.parent / "docs" / "CorpSec-Bank-Sumut-Presentasi.pptx"


def set_slide_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, top=Inches(0), height=Inches(0.08)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), top, Inches(13.333), height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.fill.background()
    return shape


def add_footer(slide, text: str):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = SLATE
    p.alignment = PP_ALIGN.RIGHT


def add_title_slide(prs, title: str, subtitle: str, extra: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_accent_bar(slide, top=Inches(0), height=Inches(0.12))

    # Decorative circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(9.5), Inches(-1), Inches(5), Inches(5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = ORANGE
    circle.fill.transparency = 0.85
    circle.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE

    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(10), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = ORANGE

    if extra:
        ex_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(10), Inches(0.8))
        tf = ex_box.text_frame
        p = tf.paragraphs[0]
        p.text = extra
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0xBC, 0xCD, 0xDC)

    add_footer(slide, "Bank Sumut — Corporate Secretary Digital Platform | PRD Discovery v0.1")


def add_section_slide(prs, section_num: str, title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY_MID)
    add_accent_bar(slide)

    num_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(2), Inches(1))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_num
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = ORANGE

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(11), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE


def add_content_slide(
    prs,
    title: str,
    bullets: list[str],
    subtitle: str = "",
    two_column: list[str] | None = None,
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_accent_bar(slide)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Orange underline
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.05), Inches(1.5), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ORANGE
    line.line.fill.background()

    y_start = 1.35
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(y_start), Inches(12), Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(13)
        p.font.color.rgb = SLATE
        y_start += 0.45

    if two_column:
        left_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(y_start), Inches(5.8), Inches(5.5)
        )
        right_box = slide.shapes.add_textbox(
            Inches(6.8), Inches(y_start), Inches(5.8), Inches(5.5)
        )
        for box, items in [(left_box, bullets), (right_box, two_column)]:
            tf = box.text_frame
            tf.word_wrap = True
            for i, item in enumerate(items):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = item
                p.font.size = Pt(14)
                p.font.color.rgb = DARK_TEXT
                p.space_after = Pt(10)
                p.level = 0
                if item.startswith("  "):
                    p.level = 1
                    p.font.size = Pt(12)
    else:
        body_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(y_start), Inches(12), Inches(5.8)
        )
        tf = body_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.font.size = Pt(15)
            p.font.color.rgb = DARK_TEXT
            p.space_after = Pt(8)
            if item.startswith("  "):
                p.level = 1
                p.font.size = Pt(13)
                p.font.color.rgb = SLATE

    add_footer(slide, "CorpSec Bank Sumut — PRD Discovery v0.1")


def add_table_slide(prs, title: str, headers: list[str], rows: list[list[str]]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_accent_bar(slide)

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY

    cols = len(headers)
    table_shape = slide.shapes.add_table(
        len(rows) + 1, cols, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.4 * (len(rows) + 2))
    )
    table = table_shape.table

    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(11)
            paragraph.font.color.rgb = WHITE

    for i, row in enumerate(rows, start=1):
        for j, value in enumerate(row):
            cell = table.cell(i, j)
            cell.text = value
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.color.rgb = DARK_TEXT

    add_footer(slide, "CorpSec Bank Sumut — PRD Discovery v0.1")


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Cover
    add_title_slide(
        prs,
        "CorpSec Bank Sumut",
        "Corporate Secretary Digital Platform",
        "Presentasi Produk — PRD Discovery v0.1\nJuli 2026",
    )

    # 2. Agenda
    add_content_slide(
        prs,
        "Agenda Presentasi",
        [
            "Latar Belakang & Tantangan Bisnis",
            "Visi Produk & Ruang Lingkup PRD",
            "Arsitektur Teknologi",
            "Modul & Fitur Utama Website",
            "Workflow Persetujuan 7 Level",
            "Role-Based Access Control (RBAC)",
            "Fitur AI Compliance Review",
            "User Journey per Role",
            "SLA, Laporan & KPI Success Metrics",
            "Desain UI & Branding Bank Sumut",
            "Kesimpulan & Roadmap",
        ],
    )

    # Section 1
    add_section_slide(prs, "01", "Latar Belakang & Tantangan")

    # 3. Background
    add_content_slide(
        prs,
        "Latar Belakang",
        [
            "Corporate Secretary Bank Sumut mengelola proses strategis: memorandum, agenda Direksi, rapat, dan compliance.",
            "Proses manual masih mengandalkan kertas, email, dan tracking terpisah — rentan terhadap keterlambatan.",
            "PRD Discovery v0.1 mendefinisikan kebutuhan digitalisasi end-to-end untuk seluruh alur CorpSec.",
            "Platform ini dibangun sebagai prototype/demo yang mengimplementasikan seluruh requirement PRD.",
        ],
        subtitle="Mengapa platform digital CorpSec dibutuhkan?",
    )

    # 4. Challenges
    add_content_slide(
        prs,
        "Tantangan vs Solusi",
        [
            "Tantangan:",
            "  • Proses approval memorandum lambat & sulit dilacak",
            "  • Dokumen fisik & penggunaan kertas masih tinggi",
            "  • Tidak ada SLA monitoring terpusat",
            "  • Review compliance manual & tidak konsisten",
            "  • Agenda Direksi & tindak lanjut rapat terfragmentasi",
        ],
        two_column=[
            "Solusi Platform:",
            "  • Workflow digital 7 level dengan notifikasi",
            "  • Upload PDF & analisa AI otomatis",
            "  • Dashboard SLA real-time (target ≥95%)",
            "  • Knowledge Base RAG untuk referensi regulasi",
            "  • Modul terintegrasi: Agenda, Rapat, Media, Laporan",
        ],
    )

    # Section 2
    add_section_slide(prs, "02", "Visi Produk & PRD")

    # 5. Vision
    add_content_slide(
        prs,
        "Visi Produk",
        [
            "Satu platform terintegrasi untuk seluruh operasional Corporate Secretary Bank Sumut.",
            "Memorandum digital dari pengusul divisi hingga keputusan final Direksi/Komisaris.",
            "Analisa AI compliance dengan referensi SMD & regulasi OJK/BI.",
            "Tanda tangan digital pada tahap approval Pimpinan Bidang & Board.",
            "Monitoring SLA, media, dan KPI dalam satu dashboard.",
        ],
        subtitle="CorpSec Bank Sumut — Digital Platform",
    )

    # 6. PRD Scope
    add_content_slide(
        prs,
        "Ruang Lingkup PRD Discovery v0.1",
        [
            "Digitalisasi agenda Direksi — CRUD jadwal & persiapan kegiatan",
            "Manajemen memorandum end-to-end — upload, review, approval, arsip",
            "Workflow persetujuan multi-level — 7 tahap dengan cabang Direksi/Komisaris",
            "AI compliance review — ekstraksi PDF, ringkasan, risk & compliance score",
            "Meeting management — rapat, notulen, tindak lanjut",
            "Media monitoring — berita Bank Sumut, regulasi, sentimen",
            "Knowledge Base — dokumen internal (SK, SE, SOP) & regulasi eksternal",
            "SLA monitoring & reporting — KPI success metrics PRD",
        ],
    )

    # Section 3
    add_section_slide(prs, "03", "Arsitektur Teknologi")

    # 7. Tech Stack
    add_table_slide(
        prs,
        "Tech Stack",
        ["Layer", "Teknologi", "Keterangan"],
        [
            ["Frontend", "Next.js 16 + React 19", "App Router, Server Components"],
            ["Bahasa", "TypeScript 5", "Type-safe development"],
            ["Styling", "Tailwind CSS v4", "Brand navy #0C2340 & orange #F58220"],
            ["Database", "SQLite + Drizzle ORM", "data/corpsec.db, auto-seed"],
            ["Grafik", "Recharts 3.9", "Dashboard & laporan statistik"],
            ["PDF", "pdf-parse", "Ekstraksi teks dari memorandum PDF"],
            ["Auth", "Cookie Session", "RBAC 7 role, middleware guard"],
            ["Ikon", "Lucide React", "UI icons konsisten"],
        ],
    )

    # 8. Architecture
    add_content_slide(
        prs,
        "Arsitektur Aplikasi",
        [
            "Struktur Next.js App Router dengan route groups per role:",
            "  • (main) — CorpSec: /dashboard, /agenda, /memorandum, /sla, /laporan",
            "  • (pengusul) — Divisi Pengusul: /pengusul, /pengusul/kirim",
            "  • Role portals — /pimpinan-bidang, /sekdireksi, /direksi, /komisaris",
            "API Routes — /api/memorandum, /api/auth, /api/agenda",
            "Middleware — canAccessRoute() proteksi akses per role",
            "Database — 11 tabel utama dengan relasi workflow & approval",
            "File storage — PDF upload di data/uploads/",
        ],
    )

    # Section 4
    add_section_slide(prs, "04", "Modul & Fitur Website")

    # 9. Modules overview
    add_table_slide(
        prs,
        "Modul Utama Website",
        ["Modul", "Route", "Fungsi"],
        [
            ["Dashboard CorpSec", "/dashboard", "Statistik, grafik, notifikasi, aksi cepat"],
            ["Memorandum", "/dashboard/memorandum", "Upload PDF, daftar, workflow, AI review"],
            ["Agenda Direksi", "/agenda", "CRUD jadwal kegiatan & persiapan"],
            ["Meeting / Rapat", "/rapat", "Rapat, notulen, tindak lanjut"],
            ["Media Monitoring", "/media", "Berita, regulasi, analisis sentimen"],
            ["Knowledge Base", "/knowledge", "Dokumen internal & regulasi (RAG)"],
            ["SLA Monitoring", "/sla", "Target vs aktual, compliance ≥95%"],
            ["Laporan & KPI", "/laporan", "Statistik & success metrics PRD"],
        ],
    )

    # 10. Dashboard
    add_content_slide(
        prs,
        "Dashboard Corporate Secretary",
        [
            "Stat cards: Agenda hari ini, Memorandum pending, SLA compliance, Notifikasi regulator",
            "Grafik interaktif: Statistik memorandum (pie), tren bulanan, tren SLA",
            "Widget operasional: Agenda hari ini, memorandum menunggu keputusan",
            "Alert: Urgensi tinggi, belum di-resume AI, tindak lanjut rapat",
            "Notifikasi keputusan Pimpinan Bidang (approve/reject) dengan badge unread",
            "Aksi cepat: Upload memorandum & buat agenda Direksi",
            "Feed berita Bank Sumut & notifikasi regulator OJK/BI",
        ],
        subtitle="Pusat kendali operasional CorpSec",
    )

    # 11. Memorandum
    add_content_slide(
        prs,
        "Manajemen Memorandum",
        [
            "Upload memorandum PDF (max 20 MB) — CorpSec atau Divisi Pengusul",
            "Metadata: No. memorandum, perihal, tanggal, divisi pengusul, urgensi",
            "Indikator baca: titik merah (belum dibaca) / oranye (sudah dibaca)",
            "Daftar memorandum dengan status workflow, AI score, aging",
            "Detail memorandum: PDF viewer, workflow stepper, panel aksi per role",
            "Tanda tangan digital pada approval Pimpinan Bidang & Board",
            "Re-upload dokumen review yang sudah diedit CorpSec",
        ],
    )

    # Section 5
    add_section_slide(prs, "05", "Workflow 7 Level")

    # 12. Workflow
    add_content_slide(
        prs,
        "Workflow Persetujuan 7 Level",
        [
            "Level 1 — Divisi Pengusul: Upload & kirim memorandum PDF",
            "Level 2 — Corporate Secretary: Analisa AI & review compliance",
            "Level 3 — Pemimpin Bidang: Approve/tolak + tanda tangan digital",
            "Level 4 — Sekretaris Direksi/Komisaris: Terima & forward ke board",
            "Level 5 — Direksi (5 anggota) / Komisaris (3 anggota): Keputusan board",
            "Level 6 — CorpSec: Finalisasi keputusan board",
            "Level 7 — Kembali ke Pengusul: Notifikasi keputusan final",
        ],
        subtitle="Alur end-to-end sesuai PRD",
        two_column=[
            "Cabang Rute:",
            "  • Rute Direksi → Sekretaris Direksi → 5 Direksi",
            "  • Rute Komisaris → Sekretaris Komisaris → 3 Komisaris",
            "",
            "Posisi Board:",
            "  • Direksi: Dirut, Dir IT, Dir Keuangan, Dir Kepatuhan, Dir Bisnis",
            "  • Komisaris: Utama, Independen, Komisaris",
            "",
            "Fitur Workflow:",
            "  • Multi-approval (semua target harus memutuskan)",
            "  • Reject dengan komentar revisi/disposisi",
            "  • Notifikasi in-app per role",
        ],
    )

    # Section 6
    add_section_slide(prs, "06", "RBAC & User Journey")

    # 13. RBAC
    add_table_slide(
        prs,
        "Role-Based Access Control — 7 Role",
        ["Role", "Home Route", "Akses Utama"],
        [
            ["Divisi Pengusul", "/pengusul", "Kirim PDF, riwayat keputusan"],
            ["Corporate Secretary", "/dashboard", "Akses penuh seluruh modul"],
            ["Pemimpin Bidang", "/pimpinan-bidang", "Review, approve/tolak + TTD"],
            ["Sekretaris Direksi", "/sekdireksi", "Terima & forward ke Direksi"],
            ["Sekretaris Komisaris", "/sekretaris-komisaris", "Terima & forward ke Komisaris"],
            ["Direksi", "/direksi", "Keputusan + tanda tangan digital"],
            ["Komisaris", "/komisaris", "Keputusan + disposisi"],
        ],
    )

    # 14. Demo accounts
    add_table_slide(
        prs,
        "Akun Demo untuk Presentasi",
        ["Role", "Username", "Password"],
        [
            ["Divisi Pengusul", "pengusul", "pengusul123"],
            ["Corporate Secretary", "corpsec", "corpsec123"],
            ["Pemimpin Bidang", "pimpinan", "pimpinan123"],
            ["Sekretaris Direksi", "sekdireksi", "sekdireksi123"],
            ["Sekretaris Komisaris", "sekkom", "sekkom123"],
            ["Direktur Utama", "dirut", "dirut123"],
            ["Direktur IT", "dir_it", "dirit123"],
            ["Komisaris Utama", "kom_utama", "kom123"],
        ],
    )

    # 15. User journeys
    add_content_slide(
        prs,
        "User Journey per Role",
        [
            "Pengusul: Login → Kirim PDF → Pantau riwayat → Terima keputusan final",
            "CorpSec: Terima memo → AI review → Kirim Pimpinan → Forward Sekretariat → Finalisasi",
            "Pimpinan Bidang: Review AI summary + PDF → Approve/Tolak + TTD digital",
            "Sekretaris: Terima memo → Assign target board → Forward ke Direksi/Komisaris",
            "Direksi/Komisaris: Review dokumen → Setujui + TTD / Tolak + disposisi",
        ],
        subtitle="Alur pengguna dari sudut pandang masing-masing role",
    )

    # Section 7
    add_section_slide(prs, "07", "Fitur AI & Compliance")

    # 16. AI
    add_content_slide(
        prs,
        "AI Compliance Review",
        [
            "Trigger otomatis saat upload PDF (autoAnalyze: true)",
            "Ekstraksi teks dari PDF menggunakan pdf-parse",
            "Analisa AI menghasilkan:",
            "  • SMD Document ID (format SMD-BSM-{tahun}-{seq})",
            "  • Ringkasan eksekutif terstruktur",
            "  • Risk Score (0–90) & Compliance Score (85–99%)",
            "  • Confidence level (82–91%)",
            "  • Referensi regulasi dari Knowledge Base (RAG keyword matching)",
            "CorpSec dapat edit ringkasan AI, download review, & re-upload",
        ],
        subtitle="Simulasi AI compliance — siap diintegrasikan dengan LLM nyata",
    )

    # 17. Knowledge & Media
    add_content_slide(
        prs,
        "Knowledge Base & Media Monitoring",
        [
            "Knowledge Base — corpus RAG untuk analisa AI:",
            "  • Dokumen internal: SK, SE, SOP Bank Sumut",
            "  • Dokumen eksternal: regulasi OJK, BI",
            "Media Monitoring — pantau reputasi & lingkungan bisnis:",
            "  • Kategori: Bank Sumut, Regulasi, Makroekonomi, Perbankan",
            "  • Analisis sentimen: Positif / Netral / Negatif",
            "  • Sumber: Bisnis Indonesia, Kontan, Reuters, dll.",
            "Notifikasi regulator OJK/BI di dashboard CorpSec",
        ],
    )

    # Section 8
    add_section_slide(prs, "08", "SLA, Laporan & KPI")

    # 18. SLA & KPI
    add_table_slide(
        prs,
        "KPI Success Metrics (PRD)",
        ["KPI", "Target PRD", "Implementasi Platform"],
        [
            ["Resume AI", "< 5 menit", "Simulasi ~2 menit otomatis"],
            ["Waktu Approval", "-50%", "Workflow digital 7 level"],
            ["Penggunaan Kertas", "-90%", "Upload & review PDF digital"],
            ["Agenda Terdokumentasi", "100%", "Modul agenda terintegrasi"],
            ["SLA Memorandum", "≥ 95%", "Monitoring real-time di /sla"],
            ["Memorandum Terlambat", "0", "Alert SLA breached di dashboard"],
        ],
    )

    # 19. Reporting
    add_content_slide(
        prs,
        "Modul SLA & Laporan",
        [
            "SLA Monitoring (/sla):",
            "  • Target compliance ≥ 95%",
            "  • Status: On Track, Berisiko, Terlambat (breached)",
            "  • Grafik tren SLA & tabel detail per memorandum",
            "Laporan & KPI (/laporan):",
            "  • Total memorandum, rapat, SLA compliance",
            "  • Grafik: tren bulanan, per status, per divisi, sentimen media",
            "  • Tabel KPI success metrics sesuai PRD",
        ],
    )

    # Section 9
    add_section_slide(prs, "09", "Desain UI & Database")

    # 20. UI Design
    add_content_slide(
        prs,
        "Desain UI — Branding Bank Sumut",
        [
            "Identitas visual corporate: Navy (#0C2340) + Orange (#F58220)",
            "Halaman login dengan hero panel gambar gedung Bank Sumut",
            "Sidebar navy gelap dengan navigasi aktif oranye",
            "Dashboard hero banner dengan gambar corporate",
            "Komponen konsisten: btn-primary, card-brand, gradient-navy",
            "Workflow stepper visual, indikator baca, badge status",
            "Responsive design — mobile sidebar dengan hamburger menu",
        ],
        subtitle="Modern, simple, corporate — tanpa warna hijau",
    )

    # 21. Database
    add_table_slide(
        prs,
        "Database Schema — Highlights",
        ["Tabel", "Fungsi"],
        [
            ["users", "Auth, role, divisi, board_position"],
            ["memorandum", "Dokumen inti + AI fields + workflow timestamps"],
            ["memorandum_approvals", "Keputusan per anggota board + TTD"],
            ["notifications", "Notifikasi in-app per user"],
            ["agenda", "Jadwal kegiatan Direksi"],
            ["meetings / meeting_followups", "Rapat & tindak lanjut"],
            ["media_articles", "Media monitoring + sentimen"],
            ["knowledge_documents", "Corpus RAG regulasi"],
            ["regulatory_notifications", "Alert OJK/BI"],
            ["sla_records", "Target vs aktual jam per memorandum"],
        ],
    )

    # Section 10
    add_section_slide(prs, "10", "Kesimpulan & Roadmap")

    # 22. Conclusion
    add_content_slide(
        prs,
        "Kesimpulan",
        [
            "Platform CorpSec Bank Sumut mengimplementasikan seluruh ruang lingkup PRD Discovery v0.1.",
            "Workflow 7 level dengan AI review, tanda tangan digital, & SLA monitoring terintegrasi.",
            "7 role dengan RBAC lengkap — dari Divisi Pengusul hingga Direksi/Komisaris.",
            "8 modul operasional: Dashboard, Memorandum, Agenda, Rapat, Media, Knowledge, SLA, Laporan.",
            "UI modern dengan branding Bank Sumut (navy & orange) dan gambar corporate.",
            "Prototype siap demo dengan data seed lengkap & 8+ akun demo.",
        ],
    )

    # 23. Roadmap
    add_content_slide(
        prs,
        "Roadmap & Next Steps",
        [
            "Integrasi LLM nyata (OpenAI/Claude) untuk analisa AI compliance",
            "Koneksi SMD (Smart Document Management) Bank Sumut",
            "Single Sign-On (SSO) dengan Active Directory Bank Sumut",
            "Audit trail & security hardening (signed session, bcrypt, rate limiting)",
            "Email/SMS notification untuk approval & SLA breach",
            "Export laporan PDF/Excel untuk reporting regulator",
            "Mobile app atau PWA untuk approval on-the-go",
        ],
        subtitle="Pengembangan lanjutan menuju production-ready",
    )

    # 24. Closing
    add_title_slide(
        prs,
        "Terima Kasih",
        "CorpSec Bank Sumut — Digital Platform",
        "Demo: http://localhost:3000/login\nBranch: cursor/banksumut-design-5a99",
    )

    return prs


def main():
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs = build_presentation()
    prs.save(str(OUTPUT))
    print(f"Presentasi berhasil dibuat: {OUTPUT}")
    print(f"Total slide: {len(prs.slides)}")


if __name__ == "__main__":
    main()
